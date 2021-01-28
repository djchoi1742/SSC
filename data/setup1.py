import argparse

parser = argparse.ArgumentParser()
setup_config = parser.add_argument_group('dataset setting')
setup_config.add_argument('--exp_name', type=str, dest='exp_name', default='exp153')
setup_config.add_argument('--data_name', type=str, dest='data_name', default='trval', help='npy name')
setup_config.add_argument('--train_excel', type=str, dest='train_excel', default='exp153_pretrain.xlsx')
setup_config.add_argument('--val_excel', type=str, dest='val_excel', default='exp153_preval.xlsx')
setup_config.add_argument('--image_prop', type=float, dest='image_prop', default=1)
setup_config.add_argument('--batch_size', type=int, dest='batch_size', default=60)
setup_config.add_argument('--w_radius', type=int, dest='w_radius', default=28)
setup_config.add_argument('--h_radius', type=int, dest='h_radius', default=21)
setup_config.add_argument('--image_size', type=int, dest='image_size', default=280)
setup_config.add_argument('--use_ssc', type=lambda x: x.title() in str(True), dest='use_ssc', default=False)
setup_config.add_argument('--aug_param', type=int, dest='aug_param', default=15)
setup_config.add_argument('--label_type', type=str, dest='label_type', default='BIN0')

parser.print_help()
config, unparsed = parser.parse_known_args()

import tensorflow as tf
import numpy as np
import pandas as pd
import os
import json
import matplotlib.pyplot as plt
import pptx
from pptx.util import Inches

import pydicom as dcm
import skimage.transform, scipy.misc
import re
import warnings

warnings.filterwarnings('ignore')

DATA_PATH = '/data/SNUBH/SSC/'  # 210.114.91.201
EXP_PATH = '/data/SNUBH/SSC/'

RAW_PATH = os.path.join(DATA_PATH, 'RAW')

train_excel_path = os.path.join(EXP_PATH, 'info', 'dataset', config.train_excel)
val_excel_path = os.path.join(EXP_PATH, 'info', 'dataset', config.val_excel)

h_w_prop = config.h_radius / config.w_radius

image_width = config.image_size
image_height = int(config.image_size * h_w_prop)

IMAGE_SIZE = [image_height, image_width]
is_flip = True
aug_param = config.aug_param

label_type = 'SSC_LABEL_' + config.label_type


def data_setting(npy_name):
    only_val = config.train_excel == config.val_excel
    npy_path = os.path.join(EXP_PATH, config.exp_name, 'npy')
    data_dir = os.path.join(npy_path, npy_name)
    excel_path = os.path.join(EXP_PATH, config.exp_name, 'excel')
    if not os.path.exists(excel_path):
        os.makedirs(excel_path)

    excel = pd.DataFrame()
    excel_val = pd.read_excel(val_excel_path)
    excel_val['DATA_TYPE'] = 'val'
    excel = excel.append(excel_val, ignore_index=True)

    if only_val is False:
        excel_train = pd.read_excel(train_excel_path)
        excel_train['DATA_TYPE'] = 'train'
        excel = excel.append(excel_train, ignore_index=True)

    select_col = ['NUMBER', 'FOLDER_NAME', 'SIDE', label_type, 'VIEW5',
                  'LX1', 'LY1', 'LX2', 'LY2',
                  'SPACING5_X', 'SPACING5_Y', 'REVERSE', 'DIRECTION',
                  'PATIENT_AGE', 'TRAUMA', 'DOMINANT_SIDE', 'VAS_MED', 'DATA_TYPE']
    excel = excel[select_col]

    def calculate_crop_coord(center_x, center_y, spacing_x, spacing_y, w_radius, h_radius,
                             side=0, direction=5, y1_adj=1.0, y2_adj=1.0):
        x1_adj, x2_adj = 1.0, 1.0
        if side == 1:  # side == LT
            x1_adj, x2_adj = 1.0, 1.0
        if direction == 0:  # flipped image
            x1_adj, x2_adj = 1.0, 1.0

        x1, y1 = int(center_x - x1_adj * w_radius / spacing_x), int(center_y - y1_adj * h_radius / spacing_y)
        x2, y2 = int(center_x + x2_adj * w_radius / spacing_x), int(center_y + y2_adj * h_radius / spacing_y)
        return [x1, y1, x2, y2]

    def add_value(excel, view_type, w_radius, h_radius):
        view_name = 'VIEW' + view_type
        spacing_x, spacing_y = 'SPACING' + view_type + '_X', 'SPACING' + view_type + '_Y'

        excel['FILES' + view_type] = excel.apply(
            lambda row: os.path.join(RAW_PATH, row['FOLDER_NAME'], row[view_name]), axis=1)

        excel['LC_COORD_X'] = excel.apply(lambda row: int((row['LX1'] + row['LX2']) / 2), axis=1)
        excel['LC_COORD_Y'] = excel.apply(lambda row: int((row['LY1'] + row['LY2']) / 2), axis=1)

        excel[view_name + '_COORD'] = excel.apply(
            lambda row: calculate_crop_coord(center_x=row['LC_COORD_X'], center_y=row['LC_COORD_Y'],
                                             spacing_x=row[spacing_x], spacing_y=row[spacing_y],
                                             w_radius=w_radius, h_radius=h_radius,
                                             y1_adj=0.8, y2_adj=1.2), axis=1)  # exp110: 1.0, 1.0

        excel['LABELS' + view_type] = excel.apply(lambda row: row[view_name + '_COORD'] + [row['SIDE']] +
                                                              [row['DIRECTION']] + [row[label_type]], axis=1)
        return excel

    def add_value_mix(excel, view_type, w_radius, h_radius):
        view_name = 'VIEW' + view_type
        spacing_x, spacing_y = 'SPACING' + view_type + '_X', 'SPACING' + view_type + '_Y'

        excel['FILES'+view_type] = excel.apply(
            lambda row: os.path.join(RAW_PATH, row['FOLDER_NAME'], row[view_name]), axis=1)

        excel['LC_COORD_Y'] = excel.apply(lambda row: int((row['L_COORD1_Y'] + row['L_COORD2_Y'])/2), axis=1)
        excel[view_name + '_COORD'] = excel.apply(
            lambda row: calculate_crop_coord(center_x=row['L_COORD2_X'], center_y=row['LC_COORD_Y'],
                                             spacing_x=row[spacing_x], spacing_y=row[spacing_y],
                                             w_radius=w_radius, h_radius=h_radius,
                                             y1_adj=1, y2_adj=1), axis=1)

        excel['LABELS'+view_type] = excel.apply(lambda row: row[view_name+'_COORD'] +
                                                            [row['SSC_LABEL_BIN0']], axis=1)
        return excel

    def combine_clinical(patient_age, trauma, dominant_side, vas_med):
        return '_'.join([str(patient_age), str(trauma), str(dominant_side), str(vas_med)])

    def files_labels_view(excel, view_type, data_type):
        files_view = excel[excel['DATA_TYPE'] == data_type]['FILES' + view_type].values
        labels_view = excel[excel['DATA_TYPE'] == data_type]['LABELS' + view_type].values
        return files_view, labels_view

    if config.use_ssc:
        val_part = excel[excel['DATA_TYPE'] == 'val']
        excel_val = add_value(val_part, '5', config.w_radius, config.h_radius)
        train_part = excel[excel['DATA_TYPE'] == 'train']
        excel_train = add_value(train_part, '5', config.w_radius, config.h_radius)
        excel = pd.concat([excel_val, excel_train], axis=0)
    else:
        excel = add_value(excel, '5', config.w_radius, config.h_radius)

    excel['CLINICAL'] = excel.apply(
        lambda row: combine_clinical(row['PATIENT_AGE'], row['TRAUMA'], row['DOMINANT_SIDE'], row['VAS_MED']),
        axis=1)

    file5_val, label5_val = files_labels_view(excel, '5', 'val')
    id_val = excel[excel['DATA_TYPE'] == 'val']['NUMBER'].values
    clv_val = excel[excel['DATA_TYPE'] == 'val']['CLINICAL'].values

    val_info = pd.DataFrame({'FILES5': pd.Series(file5_val), 'LABELS5': pd.Series(label5_val),
                             'CLINICAL': pd.Series(clv_val), 'ID': pd.Series(id_val)})

    val_size = len(val_info)

    val_info.to_csv(os.path.join(excel_path, npy_name + '_val.csv'))

    val_log = {
        'VAL_FILE': config.val_excel,
        'W_RADIUS': config.w_radius,
        'H_RADIUS': config.h_radius,
        'DESCRIPTION': 'pre-training data set: validation set',
        'LABEL_TYPE': config.label_type,
        'SIZE': val_size
    }

    with open(os.path.join(excel_path, npy_name+'_val.info'), 'w') as f:
        f.write(json.dumps(val_log, indent=4, sort_keys=True))
        f.close()

    only_val = config.train_excel == config.val_excel
    # print('only val: ', only_val)

    if only_val:
        data_set = DataSettingV1(data_dir=data_dir, batch_size=config.batch_size, only_val=only_val,
                                 image_size=config.image_size, image_height=image_height, is_flip=is_flip,
                                 file5_val=file5_val, label5_val=label5_val,
                                 id_val=id_val, clv_val=clv_val)
    else:
        file5_train, label5_train = files_labels_view(excel, '5', 'train')
        id_train = excel[excel['DATA_TYPE'] == 'train']['NUMBER'].values
        clv_train = excel[excel['DATA_TYPE'] == 'train']['CLINICAL'].values

        train_info = pd.DataFrame({'FILES5': pd.Series(file5_train), 'LABELS5': pd.Series(label5_train),
                                   'CLINICAL': pd.Series(clv_train), 'ID': pd.Series(id_train)})
        train_size = len(train_info)
        train_info.to_csv(os.path.join(excel_path, npy_name + '_train.csv'))

        train_log = {
            'TRAIN_FILE': config.train_excel,
            'W_RADIUS': config.w_radius,
            'H_RADIUS': config.h_radius,
            'DESCRIPTION': 'pre-training data set: training set',
            'LABEL_TYPE': config.label_type,
            'SIZE': train_size
        }

        with open(os.path.join(excel_path, npy_name + '_train.info'), 'w') as f:
            f.write(json.dumps(train_log, indent=4, sort_keys=True))
            f.close()

        data_set = DataSettingV1(data_dir=data_dir, batch_size=config.batch_size, only_val=only_val,
                                 image_size=config.image_size, image_height=image_height, is_flip=is_flip,
                                 file5_val=file5_val, label5_val=label5_val,
                                 id_val=id_val, clv_val=clv_val,
                                 file5_train=file5_train, label5_train=label5_train,
                                 id_train=id_train, clv_train=clv_train
                                 )
    return data_set


class DataSettingV1:
    def __init__(self, data_dir, batch_size, only_val, image_size, image_height, replicate=False,
                 rep_label=1, rep_count=2, aug_param=20, is_flip=True, **kwargs):
        if not os.path.exists(data_dir):
            if 'clv_train' in kwargs:
                train_x5, train_y5 = kwargs['file5_train'], kwargs['label5_train']
                train_id, train_clv = kwargs['id_train'], kwargs['clv_train']
                
            if 'clv_val' in kwargs:
                val_x5, val_y5 = kwargs['file5_val'], kwargs['label5_val']
                val_id, val_clv = kwargs['id_val'], kwargs['clv_val']

            else:
                raise AssertionError('files or labels must be provided. please check npy file.')

            data_root = os.path.split(data_dir)[0]
            if not os.path.exists(data_root): os.makedirs(data_root)

            if only_val:
                np.save(data_dir, {'val_x5': val_x5, 'val_y5': val_y5,
                                   'val_id': val_id, 'val_clv': val_clv
                                   })
            else:
                np.save(data_dir, {'train_x5': train_x5, 'train_y5': train_y5,
                                   'train_id': train_id, 'train_clv': train_clv,
                                   'val_x5': val_x5, 'val_y5': val_y5,
                                   'val_id': val_id, 'val_clv': val_clv
                                   })
        else:
            pre_built = np.load(data_dir).item()

            val_x5, val_y5 = pre_built['val_x5'], pre_built['val_y5']
            val_id, val_clv = pre_built['val_id'], pre_built['val_clv']
            self.data_length = len(val_id)

            if only_val is False:
                train_x5, train_y5 = pre_built['train_x5'], pre_built['train_y5']
                train_id, train_clv = pre_built['train_id'], pre_built['train_clv']
                self.data_length = len(train_id) + len(val_id)

        self.val = self.SubDataSetting((val_x5, val_y5, val_id, val_clv), batch_size=batch_size,
                                       image_size=image_size, image_height=image_height, is_flip=is_flip,
                                       aug_param=aug_param, shuffle=False, augmentation=False)
        self.id_val = val_id

        if only_val is False:
            if replicate:
                def replicate_data(rep_x, rep_y, rep_id, rep_clv, label, rep_num):
                    index = np.asarray([v[-1] for v in rep_y])

                    rep_label_x = rep_x[index == label]
                    rep_label_y = rep_y[index == label]
                    rep_label_id = rep_id[index == label]
                    rep_label_clv = rep_clv[index == label]

                    for i in range(0, rep_num):
                        rep_x = np.concatenate([rep_x, rep_label_x], axis=0)
                        rep_y = np.concatenate([rep_y, rep_label_y], axis=0)
                        rep_id = np.concatenate([rep_id, rep_label_id], axis=0)
                        rep_clv = np.concatenate([rep_clv, rep_label_clv], axis=0)

                    rep_out_x, rep_out_y, rep_out_id, rep_out_clv = rep_x, rep_y, rep_id, rep_clv

                    return rep_out_x, rep_out_y, rep_out_id, rep_out_clv

                train_x5, train_y5, train_id, train_clv = \
                    replicate_data(train_x5, train_y5, train_id, train_clv, rep_label, rep_count)

                print(len(train_x5))

            np.random.seed(20191224)
            p = np.random.permutation(len(train_x5))
            train_x5, train_y5 = train_x5[p], train_y5[p]
            train_id, train_clv = train_id[p], train_clv[p]

            self.train = self.SubDataSetting((train_x5, train_y5, train_id, train_clv), batch_size=batch_size,
                                             image_size=image_size, image_height=image_height, is_flip=is_flip,
                                             aug_param=aug_param, shuffle=True, augmentation=True)

            self.id_train = train_id

    class SubDataSetting:
        def __init__(self, files_n_labels, num_epochs=1, batch_size=1,
                     image_height=image_height, image_size=image_width, is_flip=is_flip, aug_param=aug_param,
                     shuffle=False, augmentation=False):

            self.file5, self.label5, self.id, self.clv = files_n_labels
            self.data_length = len(self.id)

            data_set = tf.data.Dataset.from_tensor_slices(tensors=
                                                          (self.file5,
                                                           [v for v in self.label5],
                                                           [v for v in self.id],
                                                           [v for v in self.clv]
                                                           ))
            if shuffle:
                data_set = data_set.shuffle(buffer_size=batch_size * 100, reshuffle_each_iteration=True)

            def dcm_read_by_ftn(file5, label5, id, clv, image_size, augmentation):
                def each_read(filename, label, augmentation):
                    dcm_info = dcm.read_file(filename.decode())
                    x1, y1, x2, y2, side, direction = label[:-1]
                    y = np.int64(label[-1])

                    if augmentation:
                        shift_x = np.random.randint(dcm_info.Columns // aug_param)  # 20
                        shift_y = np.random.randint(dcm_info.Rows // aug_param)  # 20

                        shift_x = -shift_x if np.random.rand() <= 0.5 else shift_x
                        shift_y = -shift_y if np.random.rand() <= 0.5 else shift_y

                        x1, y1 = x1 - shift_x, y1 - shift_y
                        x2, y2 = x2 - shift_x, y2 - shift_y

                    image = dcm_info.pixel_array

                    # recover image if inverted
                    if str(dcm_info[0x28, 0x04].value) == 'MONOCHROME1':
                        white_image = np.full_like(image, np.max(image), image.dtype)
                        image = np.subtract(white_image, image)

                    image = image[max(0, y1):min(dcm_info.Rows, y2), max(0, x1):min(dcm_info.Columns, x2)]

                    if is_flip:
                        if side == 1:  # side == LT
                            image = np.fliplr(image)
                        if direction == 0:  # flipped image
                            image = np.fliplr(image)

                    image_shape = [image_height, image_size]
                    image = np.expand_dims(skimage.transform.resize(image, image_shape, preserve_range=True), axis=-1)

                    image = (image - np.mean(image)) / np.std(image)  # normalization
                    image = image.astype(np.float32)

                    return image, y, side

                clv_split = re.split('_', clv.decode())
                clv_split = [np.float32(i) for i in clv_split]

                age, tm, dm, vas = clv_split[:]

                f5, l5, side = each_read(file5, label5, augmentation)

                idv = id.decode()

                return f5, l5, idv, age, tm, dm, vas, side

            data_set = data_set.map(num_parallel_calls=8,
                                    map_func=lambda file5, label5, id, clv:
                                    tuple(tf.py_func(func=dcm_read_by_ftn,
                                                     inp=[file5, label5, id, clv, image_size, augmentation],
                                                     Tout=[tf.float32, tf.int64, tf.string,
                                                           tf.float32, tf.float32, tf.float32, tf.float32,
                                                           tf.int32]
                                                     )))

            if num_epochs == 0:
                data_set = data_set.repeat(count=num_epochs)
                data_set = data_set.batch(batch_size=batch_size, drop_remainder=True)
            else:
                data_set = data_set.batch(batch_size)
                data_set = data_set.repeat(count=num_epochs)

            iterator = data_set.make_initializable_iterator()

            self.init_op = iterator.initializer
            self.next_batch = iterator.get_next()


def get_info(data_path, info_name, key):
    with open(os.path.join(data_path, info_name+'.info'), 'r') as f:
        info = json.load(f)
        return info[key]


def test_data_loader(data_set):
    check_path = os.path.join(EXP_PATH, config.exp_name, 'view')
    if not os.path.exists(check_path):
        os.makedirs(check_path)

    with tf.Session() as sess:
        sess.run(data_set.val.init_op)
        num_examples, next_batch = data_set.val.data_length, data_set.val.next_batch

        count = 0
        num_iter = int(np.ceil(float(num_examples) / config.batch_size))
        print('num_iter: ', num_iter)

        prs = pptx.Presentation()
        prs.slide_width = Inches(10 * 2)
        prs.slide_height = Inches(6 * 2)

        while count < num_iter:
            img, lbl, name, age, tm, dm, vas, _ = sess.run(data_set.val.next_batch)
            # img, lbl, name, age, sxf, sxm, vas, tm0, tm1, tm2, dm0, dm1, dm2 = \
            #    sess.run(data_set.val.next_batch)

            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)

            show_images(img, name, num_rows=6, num_cols=10)
            fig_name = '_'.join([config.exp_name, config.data_name, '%03d' % count]) + '.png'

            fig_path = os.path.join(check_path, fig_name)
            plt.savefig(fig_path, bbox_inches='tight')
            slide.shapes.add_picture(fig_path, Inches(0), Inches(0), width=Inches(10 * 2))
            os.remove(fig_path)
            count += 1

            if count % 10 == 0:
                print(count)

    ppt_name = os.path.join(check_path, '_'.join([config.exp_name, config.data_name]) + '.pptx')
    prs.save(ppt_name)
    print('Saved: ', ppt_name)


def show_images(images, names, num_rows=6, num_cols=10, fig_size=(10*2, 6*2)):
    plt.figure(figsize=fig_size)
    num_figs = images.shape[0]  # num_rows * num_cols
    # num_chars = 5  # num of chars to show in names

    for j in range(num_figs):
        plt.subplot(num_rows, num_cols, j + 1)
        plt.imshow(np.squeeze(images[j]), cmap='gray')
        plt.axis('off')
        img_name = os.path.basename(names[j])
        plt.title(str(img_name.decode('utf-8')), fontsize=8, color='blue')


if __name__ == '__main__':
    os.environ['CUDA_VISIBLE_DEVICES'] = '1'

    data_name = config.data_name
    d_set = data_setting(data_name)
    test_data_loader(d_set)




