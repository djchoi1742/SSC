import argparse

license = """
Copyright ⓒ Dongjun Choi, Kyong Joon Lee
Department of Radiology at Seoul National University Bundang Hospital. \n
If you have any question, please email us for assistance: chzze4582@gmail.com \n """
parser = argparse.ArgumentParser(formatter_class=argparse.RawDescriptionHelpFormatter, \
                                 description='', epilog=license, add_help=False)

network_config = parser.add_argument_group('network setting (must be provided)')

network_config.add_argument('--data_path', type=str, dest='data_path', default='/data/SNUBH/SSC/')
network_config.add_argument('--pre_exp', type=str, dest='pre_exp', default='exp153')
network_config.add_argument('--pre_model1', type=str, dest='pre_model', default='Model58')
# network_config.add_argument('--pre_serial1', type=int, dest='pre_serial', default=5)

network_config.add_argument('--exp_name', type=str, dest='exp_name', default='exp151')
network_config.add_argument('--model_name', type=str, dest='model_name', default='Model63')
network_config.add_argument('--batch_size', type=int, dest='batch_size', default=8)

network_config.add_argument('--trial_serial', type=str, dest='trial_serial', default='12,17,18,15')
network_config.add_argument('--esb_serial', type=int, dest='esb_serial', default=2)
network_config.add_argument('--npy_name', type=str, dest='npy_name', default='trval.npy')

network_config.add_argument('--num_weight', type=int, dest='num_weight', default=3)  # only use validation
network_config.add_argument('--image_size', type=int, dest='image_size', default=280)
network_config.add_argument('--theta', type=float, dest='theta', default=0.5)

parser.print_help()
config, unparsed = parser.parse_known_args()


import sys, os
sys.path.append('/home/chzze/bitbucket/SSC')
os.environ['TF_CPP_MIN_LOG_LEVEL'] = '2'  # ignore SSE instruction warning on tensorflow

import tensorflow as tf
import numpy as np
import sklearn.metrics  # roc curve
import matplotlib.pyplot as plt
import pandas as pd
import json, re, pptx
from pptx.util import Inches
from data.setup2 import DataSettingV1, get_info

import models.model2 as model_y
from tf_utils.result import *


try:
    trial_serial_str = '%03d' % int(config.trial_serial)
    trial_serial_list = []
    is_ensemble = False
except:
    is_ensemble = True
    trial_serial_list = re.split(',', config.trial_serial)
    trial_serial_str = 'e' + '%03d' % int(config.esb_serial)

try:
    previous_path = os.path.join(config.data_path, config.pre_exp, config.pre_model,
                                 'result-%03d' % int(trial_serial_list[0]))
    # growth_k = model_y.get_train_info(previous_path, 'GROWTH_K')
    # block_rep = model_y.get_train_info(previous_path, 'BLOCK_REP')
    # dense_type = model_y.get_train_info(previous_path, 'DENSE_TYPE')
    # last_kp = model_y.get_train_info(previous_path, 'LAST_KP')
    # use_se = model_y.get_train_info(previous_path, 'USE_SE')
    # k_p = model_y.get_train_info(previous_path, 'K_P')
except:
    pass
    # growth_k, block_rep = 32, '1,1,1,1'
    # dense_type, last_kp, use_se, k_p = 1, 1, False, '1,1,1,1'

result_path = os.path.join(config.data_path, config.exp_name, config.model_name, # config.pre_model,
                           'result-%s' % trial_serial_str)
ppt_path = os.path.join(config.data_path, config.exp_name, config.model_name, # config.pre_model,
                        'ppt-%s' % trial_serial_str)

cam_path = os.path.join(config.data_path, 'cam')


print('result_path: ', result_path)
if not os.path.exists(result_path): os.makedirs(result_path)
if not os.path.exists(ppt_path): os.makedirs(ppt_path)


ckpt_path = os.path.join(result_path, 'ckpt')
npy_path = os.path.join(config.data_path, config.exp_name, 'npy')
excel_path = os.path.join(config.data_path, config.exp_name, 'excel')

data_name, _ = os.path.splitext(config.npy_name)

h_radius = get_info(excel_path, data_name+'_val', 'H_RADIUS')
w_radius = get_info(excel_path, data_name+'_val', 'W_RADIUS')

h_w_prop = h_radius / w_radius
image_height = int(config.image_size * h_w_prop)

print(npy_path + '/' + config.npy_name)
# print('h_radius: ', h_radius, '  ', 'w_radius: ', w_radius)

data_set = DataSettingV1(data_dir=os.path.join(npy_path, config.npy_name),
                         batch_size=config.batch_size, only_val=True,
                         image_size=config.image_size, image_height=image_height)

infer_name = 'Inference'+config.model_name


model = getattr(model_y, infer_name)(trainable=False, theta=config.theta,
                                     pre_model='Inference'+config.pre_model,
                                     image_size=config.image_size, image_height=image_height)
                                     # growth_k=growth_k, block_rep=block_rep,
                                     # dense_type=dense_type, last_kp=last_kp,
                                     # use_se=use_se, k_p=k_p)


def restore_weight(data_path, exp_name, model_name, pre_model, trial_serial, num_weight):
    # weight_auc_path = os.path.join(data_path, exp_name, model_name, pre_model, 'result-%03d' % trial_serial)
    weight_auc_path = os.path.join(data_path, exp_name, model_name, 'result-%03d' % trial_serial)
    weight_auc_csv = pd.read_csv(os.path.join(weight_auc_path, '_'.join([exp_name, model_name,
                                                                         '%03d' % trial_serial])+'.csv'))
    weight_auc_csv = weight_auc_csv.sort_values('AUC', ascending=False)
    all_ckpt_paths = list(weight_auc_csv['WEIGHT_PATH'][0:int(num_weight)])
    return all_ckpt_paths


def validation():
    sess_config = tf.ConfigProto(log_device_placement=False)
    sess_config.gpu_options.allow_growth = True

    saver = tf.train.Saver()
    num_examples = len(data_set.id_val)

    if is_ensemble is False:
        all_ckpt_paths = restore_weight(config.data_path, config.exp_name, config.model_name, config.pre_model,
                                        int(config.trial_serial), config.num_weight)
    else:
        all_ckpt_paths = []
        for idx in trial_serial_list:
            each_ckpt_paths = restore_weight(config.data_path, config.exp_name, config.model_name, config.pre_model,
                                             int(idx), config.num_weight)
            all_ckpt_paths = all_ckpt_paths + each_ckpt_paths

    num_ckpt = len(all_ckpt_paths)
    print('num_ckpt: ', num_ckpt)

    info_log = {
        'esb_ckpt': all_ckpt_paths
    }

    with open(os.path.join(result_path, 'esb.info'), 'w') as f:
        f.write(json.dumps(info_log, indent=4, sort_keys=True))
        f.close()

    imgs = np.zeros([num_examples, model.img_h, model.img_w, model.img_c])
    cams = np.zeros([num_ckpt, num_examples, model.img_h, model.img_w, model.img_c])
    cams_raw = np.zeros([num_ckpt, num_examples, 7, 9, model.img_c])

    lbls = np.zeros([num_examples, ], dtype=np.int32)
    lgts = np.zeros([num_ckpt, num_examples, 1])
    probs = np.zeros([num_ckpt, num_examples, 1])

    val_x, val_y = None, None
    for ckpt_idx, ckpt_path in enumerate(all_ckpt_paths):
        print('Restoring: ' + ckpt_path)

        sess = tf.Session(config=sess_config)
        saver.restore(sess, ckpt_path)

        sess.run(data_set.val.init_op)
        val_x, val_y = [], []

        num_iter = int(np.ceil(float(num_examples) / config.batch_size))
        step = 0

        while step < num_iter:
            sys.stdout.write('Evaluation [{0}/{1}]\r'.format(len(val_y) // config.batch_size,
                             -(-data_set.val.data_length // config.batch_size)))

            img, lbl, name, age, tm, dm, vas, _ = sess.run(data_set.val.next_batch)


            feed_dict = {model.images: img, model.labels: lbl,
                         model.ages: age, model.vas: vas, model.tm: tm, model.dm: dm, model.is_training: False}

            # feed_dict = {model.images: img, model.labels: lbl, model.is_training: False}

            val_loss, val_lgt, val_prob, val_acc = \
                sess.run([model.loss, model.logits, model.prob, model.accuracy], feed_dict=feed_dict)

            if np.ndim(val_prob) == 1:
                val_prob = np.expand_dims(val_prob, axis=-1)

            val_x.extend(val_prob)
            val_y.extend(lbl)

            cam, cam_raw = sess.run([model.local, model.local1], feed_dict=feed_dict)
            
            # cam_1, cam_2, cam_3 = cam[0], cam[1], cam[2]
            # import pdb; pdb.set_trace()
            # cam = sess.run(model.local, feed_dict=feed_dict)

            cams[ckpt_idx, step * config.batch_size:step * config.batch_size + len(lbl)] = cam
            cams_raw[ckpt_idx, step * config.batch_size:step * config.batch_size + len(lbl)] = cam_raw
            probs[ckpt_idx, step * config.batch_size:step * config.batch_size + len(lbl)] = val_prob
            lgts[ckpt_idx, step * config.batch_size:step * config.batch_size + len(lbl)] = val_lgt

            if ckpt_idx == 0:
                imgs[step * config.batch_size:step * config.batch_size + len(lbl)] = img
                lbls[step * config.batch_size:step * config.batch_size + len(lbl)] = lbl

            step += 1

        sess.close()

    probs, lgts, cams = np.mean(probs, axis=0), np.mean(lgts, axis=0), np.mean(cams, axis=0)
    # cams_raw = np.mean(cams_raw, axis=0)
    id_test = data_set.id_val

    prob_1, lgts_1 = np.squeeze(np.array(probs)), np.squeeze(np.array(lgts))

    result_csv = pd.DataFrame({'NUMBER': id_test, 'PROB': prob_1, 'LOGIT': lgts_1, 'LABEL': np.array(lbls)})
    result_name = '_'.join([config.model_name, config.npy_name, trial_serial_str, '%03d' % config.num_weight])+'.csv'
    result_csv.to_csv(os.path.join(result_path, result_name), index=False)

    fpr, tpr, _ = sklearn.metrics.roc_curve(val_y, prob_1, drop_intermediate=False)
    val_auc = sklearn.metrics.auc(fpr, tpr)

    print('Validation AUC: ', val_auc)
    print('Validation Complete...\n')

    prs = pptx.Presentation()
    prs.slide_width, prs.slide_height = Inches(8*2), Inches(5*2)

    plt_batch = 20
    plt_step = 0
    plt_iter, plt_examples = int(np.ceil(num_examples / plt_batch)), num_examples

    while plt_step < plt_iter:

        if plt_examples >= plt_batch:
            len_batch = plt_batch
        else:
            len_batch = plt_examples

        images_batch = imgs[plt_step * plt_batch:plt_step * plt_batch + len_batch]
        labels_batch = lbls[plt_step * plt_batch:plt_step * plt_batch + len_batch]
        names_batch = id_test[plt_step * plt_batch:plt_step * plt_batch + len_batch]

        probs_batch = probs[plt_step * plt_batch:plt_step * plt_batch + len_batch]
        cams_batch = cams[plt_step * plt_batch:plt_step * plt_batch + len_batch]

        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)

        show_cam(cams_batch, probs_batch, images_batch, labels_batch, names_batch, 'LABEL')
        fig_name = '_'.join([config.exp_name, config.model_name, config.npy_name, trial_serial_str,
                             '%03d' % plt_step]) + '.png'
        fig_path = os.path.join(cam_path, fig_name)
        plt.savefig(fig_path, bbox_inches='tight')
        slide.shapes.add_picture(fig_path, Inches(0), Inches(0), width=Inches(8 * 2))
        # os.remove(fig_path)
        plt_step += 1
        plt_examples -= plt_batch

    print('plt_examples check: ', plt_examples)
    ppt_name = os.path.join(ppt_path, '_'.join([config.exp_name, config.model_name, config.npy_name, trial_serial_str,
                                                '%03d' % config.num_weight]) + '.pptx')
    prs.save(ppt_name)
    print('Saved: ', ppt_name)


def show_cam(cams, probs, images, labels, names, side_label, num_rows=5, num_cols=8, fig_size=(8*2, 5*2)):
    batch_size = cams.shape[0]
    fig, ax = plt.subplots(num_rows, num_cols, figsize=fig_size)
    axoff_fun = np.vectorize(lambda ax: ax.axis('off'))
    axoff_fun(ax)

    for i in range(batch_size):
        prob = '%.2f' % probs[i]
        lbl = int(labels[i])
        show_image = np.squeeze(images[i])
        cam = np.squeeze(cams[i])
        img_row, img_col = int(i % num_rows), int(i / num_rows) * 2

        ori_title = ' '.join([names[i], side_label + ': '+str(lbl)])
        cam_title = side_label+' Pred: '+str(prob)

        ax[img_row, img_col].imshow(show_image, cmap='bone')
        ax[img_row, img_col+1].imshow(show_image, cmap='bone')
        ax[img_row, img_col+1].imshow(cam, cmap=plt.cm.seismic, alpha=0.5, interpolation='nearest')

        if (lbl == 0 and probs[i] < 0.5) or (lbl == 1 and probs[i] >= 0.5):
            txt_color = 'blue'
        else:
            txt_color = 'red'
        ax[img_row, img_col].set_title(ori_title, fontsize=7, color=txt_color)
        ax[img_row, img_col+1].set_title(cam_title, fontsize=7, color=txt_color)


if __name__ == '__main__':
    validation()

