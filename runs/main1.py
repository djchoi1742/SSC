import argparse

license = """
Copyright ⓒ Dongjun Choi, Kyong Joon Lee
Department of Radiology at Seoul National University Bundang Hospital. \n
If you have any question, please email us for assistance: chzze4582@gmail.com \n """
parser = argparse.ArgumentParser(formatter_class=argparse.RawDescriptionHelpFormatter, \
                                 description='', epilog=license, add_help=False)

network_config = parser.add_argument_group('network setting (must be provided)')

network_config.add_argument('--data_path', type=str, dest='data_path', default='/data/SNUBH/SSC/')
network_config.add_argument('--exp_name', type=str, dest='exp_name', default='exp153')
network_config.add_argument('--model_name', type=str, dest='model_name', default='Model58')
network_config.add_argument('--train', type=lambda x: x.title() in str(True), dest='train', default=False)
network_config.add_argument('--batch_size', type=int, dest='batch_size', default=8)  # set multiples of 3
network_config.add_argument('--num_epoch', type=int, dest='num_epoch', default=80)  # infinite loop
network_config.add_argument('--trial_serial', type=int, dest='trial_serial', default=30)
network_config.add_argument('--npy_name', type=str, dest='npy_name', default='trval.npy')
network_config.add_argument('--save_epoch', type=int, dest='save_epoch', default=1)  # only use training
network_config.add_argument('--max_keep', type=int, dest='max_keep', default=2)  # only use training
network_config.add_argument('--num_weight', type=int, dest='num_weight', default=1)  # only use validation
network_config.add_argument('--image_size', type=int, dest='image_size', default=280)
network_config.add_argument('--val_measure', type=str, dest='val_measure', default='auc')
network_config.add_argument('--pre_weight', type=str, dest='pre_weight', default='auc')  # only use validation
network_config.add_argument('--select_weight', type=int, dest='select_weight', default=0)  # only use validation

network_config.add_argument('--learning_rate', type=float, dest='learning_rate', default=0.01)
network_config.add_argument('--decay_steps', type=int, dest='decay_steps', default=7500)
network_config.add_argument('--decay_rate', type=float, dest='decay_rate', default=0.94)
network_config.add_argument('--decay', type=float, dest='decay', default=0.9)
network_config.add_argument('--replicate', type=lambda x: x.title() in str(True), dest='replicate', default=False)
network_config.add_argument('--rep_label', type=int, dest='rep_label', default=1)
network_config.add_argument('--rep_count', type=int, dest='rep_count', default=1)
network_config.add_argument('--aug_param', type=int, dest='aug_param', default=15)

# network_config.add_argument('--use_focal', type=lambda x: x.title() in str(True), dest='use_focal', default=False)
network_config.add_argument('--alpha', type=float, dest='alpha', default=0.3)
network_config.add_argument('--gamma', type=float, dest='gamma', default=2)
network_config.add_argument('--use_dropout', type=lambda x: x.title() in str(True), dest='use_dropout', default=False)
network_config.add_argument('--use_se', type=lambda x: x.title() in str(True), dest='use_se', default=False)
network_config.add_argument('--growth_k', type=int, dest='growth_k', default=32)
network_config.add_argument('--theta', type=float, dest='theta', default=0.5)
network_config.add_argument('--block_rep', type=str, dest='block_rep', default='2,2,2')

parser.print_help()
config, unparsed = parser.parse_known_args()

import sys, os
sys.path.append('/home/chzze/bitbucket/SSC')

import tensorflow as tf
import numpy as np
import sklearn.metrics  # roc curve
import matplotlib.pyplot as plt
import pandas as pd
import json, pptx, re
from pptx.util import Inches

from data.setup1 import DataSettingV1, get_info
import models.model1 as model_z
from tf_utils.tboard import TensorBoard
from tf_utils.result import *

trial_serial_str = '%03d' % config.trial_serial
log_path = os.path.join(config.data_path, config.exp_name, config.model_name, 'logs-%s' % trial_serial_str)
result_path = os.path.join(config.data_path, config.exp_name, config.model_name, 'result-%s' % trial_serial_str)
ckpt_path = os.path.join(result_path, 'ckpt')
npy_path = os.path.join(config.data_path, config.exp_name, 'npy')
excel_path = os.path.join(config.data_path, config.exp_name, 'excel')

cam_path = os.path.join(config.data_path, 'cam')
ppt_path = os.path.join(config.data_path, config.exp_name, config.model_name, 'ppt-%s' % trial_serial_str)

if not os.path.exists(result_path):
    os.makedirs(result_path)
if not os.path.exists(cam_path):
    os.makedirs(cam_path)
if not os.path.exists(ppt_path):
    os.makedirs(ppt_path)

print(npy_path + '/' + config.npy_name)

data_name, _ = os.path.splitext(config.npy_name)

if config.train:
    h_radius = get_info(excel_path, data_name+'_train', 'H_RADIUS')
    w_radius = get_info(excel_path, data_name+'_train', 'W_RADIUS')
else:
    h_radius = get_info(excel_path, data_name+'_val', 'H_RADIUS')
    w_radius = get_info(excel_path, data_name+'_val', 'W_RADIUS')

h_w_prop = h_radius / w_radius
image_height = int(config.image_size * h_w_prop)

data_set = DataSettingV1(data_dir=os.path.join(npy_path, config.npy_name),
                         batch_size=config.batch_size, only_val=bool(1 - config.train),
                         image_size=config.image_size, image_height=image_height,
                         replicate=config.replicate, rep_label=config.rep_label,
                         rep_count=config.rep_count, aug_param=config.aug_param, is_flip=True)

infer_name = 'Inference' + config.model_name

if config.train:
    growth_k, theta = config.growth_k, config.theta
    block_rep, use_se = config.block_rep, config.use_se

else:
    growth_k = model_z.get_train_info(result_path, 'GROWTH_K')
    theta = model_z.get_train_info(result_path, 'THETA')
    block_rep = model_z.get_train_info(result_path, 'BLOCK_REP')
    use_se = model_z.get_train_info(result_path, 'USE_SE')



model = getattr(model_z, infer_name)(trainable=config.train, growth_k=growth_k, theta=theta,
                                     image_size=config.image_size, image_height=image_height,
                                     learning_rate=config.learning_rate, decay_steps=config.decay_steps,
                                     decay_rate=config.decay_rate, decay=config.decay, use_se=use_se,
                                     block_rep=block_rep, alpha=config.alpha, gamma=config.gamma)


def training():
    info_log = {
        'NPY_NAME': config.npy_name,
        'IMAGE_SIZE': config.image_size,
        'VAL_MEASURE': config.val_measure,
        'BATCH_SIZE': config.batch_size,
        'LEARNING_RATE': config.learning_rate,
        'DECAY_STEPS': config.decay_steps,
        'DECAY_RATE': config.decay_rate,
        'REPLICATE': config.replicate,
        'REP_LABEL': config.rep_label,
        'REP_COUNT': config.rep_count,
        'AUG_PARAM': config.aug_param,
        'NUM_EPOCH': config.num_epoch,
        'SAVE_EPOCH': config.save_epoch,
        'FOCAL_LOSS_ALPHA': config.alpha,
        'FOCAL_LOSS_GAMMA': config.gamma,
        'USE_DROPOUT': config.use_dropout,
        'USE_SE': config.use_se,
        'GROWTH_K': config.growth_k,
        'THETA': config.theta,
        'BLOCK_REP': config.block_rep
    }

    with open(os.path.join(result_path, '.info'), 'w') as f:
        f.write(json.dumps(info_log, indent=4, sort_keys=True))
        f.close()

    sess_config = tf.ConfigProto(log_device_placement=False, allow_soft_placement=True)
    sess_config.gpu_options.allow_growth = True
    sess = tf.Session(config=sess_config)

    t_board = TensorBoard(log_dir=log_path, overwrite=True)
    loss_rec = tf.get_variable(name='Loss', shape=[], trainable=False, initializer=tf.zeros_initializer(),
                               collections=['scalar'])
    auc_rec = tf.get_variable(name='AUC', shape=[], trainable=False, initializer=tf.zeros_initializer(),
                              collections=['scalar'])
    accuracy_rec = tf.get_variable(name='Accuracy', shape=[], trainable=False, initializer=tf.zeros_initializer(),
                                   collections=['scalar'])

    t_board.init_scalar(collections=['scalar'])
    t_board.init_images(collections=['images'], num_outputs=4)

    sess.run(tf.global_variables_initializer())
    saver = tf.train.Saver(max_to_keep=config.max_keep)

    result_name = '_'.join([config.exp_name, config.model_name, trial_serial_str]) + '.csv'
    df_csv = pd.DataFrame({'WEIGHT_PATH': pd.Series(), 'AUC': pd.Series(), 'LOSS': pd.Series()})

    current_step, current_epoch = None, None
    perf_per_epoch, max_perf_per_epoch, max_current_step = [], [], []  # measure: AUC or Loss

    try:
        while True:  # increment of training epochs
            sess.run([data_set.train.init_op, data_set.val.init_op])
            train_loss_batch = []
            train_x, train_y = [], []
            train_acc_batch = []

            train_length = data_set.train.data_length
            num_iter_train = int(np.ceil(float(train_length) / config.batch_size))
            train_step = 0

            feed_dict = {}
            while train_step < num_iter_train:

                img, lbl, name, age, tm, dm, vas, _ = sess.run(data_set.train.next_batch)
                feed_dict = {model.images: img, model.labels: lbl, model.is_training: True}

                _, train_loss, train_prob, train_acc = \
                    sess.run([model.train, model.loss, model.prob, model.accuracy], feed_dict=feed_dict)

                current_step, current_epoch = sess.run([tf.train.get_global_step(), model.global_epoch])
                sys.stdout.write('Step: {0:>4d} ({1})\r'.format(current_step, current_epoch))

                train_x.extend(train_prob)
                train_y.extend(lbl)
                train_acc_batch.append(train_acc)
                train_loss_batch.append(train_loss)

                train_step += 1

            sess.run(tf.assign_add(model.global_epoch, 1))

            fpr, tpr, _ = sklearn.metrics.roc_curve(train_y, train_x, drop_intermediate=False)
            train_auc = sklearn.metrics.auc(fpr, tpr)

            feed_dict.update({loss_rec: np.mean(train_loss_batch), auc_rec: train_auc,
                              accuracy_rec: np.mean(train_acc_batch)})

            t_board.add_summary(sess=sess, feed_dict=feed_dict, log_type='train')

            # validation with roc-auc
            val_length = data_set.val.data_length
            num_iter_val = int(np.ceil(float(val_length) / config.batch_size))
            val_step = 0

            val_loss_batch = []
            val_x, val_y = [], []
            val_acc_batch = []

            feed_dict = {}
            while val_step < num_iter_val:
                sys.stdout.write('Evaluation [{0}/{1}]\r'.format(len(val_loss_batch) + 1,
                                                                 -(-data_set.val.data_length // config.batch_size)))

                img, lbl, name, age, tm, dm, vas, _ = sess.run(data_set.val.next_batch)

                feed_dict = {model.images: img, model.labels: lbl, model.is_training: False}

                val_loss, val_prob, val_acc = sess.run([model.loss, model.prob, model.accuracy], feed_dict=feed_dict)

                val_x.extend(val_prob)
                val_y.extend(lbl)
                val_acc_batch.append(val_acc)
                val_loss_batch.append(val_loss)

                val_step += 1

            fpr, tpr, _ = sklearn.metrics.roc_curve(val_y, val_x, drop_intermediate=False)
            val_auc = sklearn.metrics.auc(fpr, tpr)
            val_loss_mean = np.mean(val_loss_batch)

            feed_dict.update({loss_rec: val_loss_mean, auc_rec: val_auc,
                              accuracy_rec: np.mean(val_acc_batch)})

            t_board.add_summary(sess=sess, feed_dict=feed_dict, log_type='val')
            t_board.display_summary(time_stamp=True)

            if config.val_measure == 'loss':
                v_measure = val_loss_mean
                if current_epoch >= config.max_keep + 1:
                    v_criteria = v_measure < max(df_csv['LOSS'].tolist())
                else:
                    v_criteria = None

            elif config.val_measure == 'auc':
                v_measure = val_auc
                if current_epoch >= config.max_keep + 1:
                    v_criteria = v_measure > min(df_csv['AUC'].tolist())
                else:
                    v_criteria = None

            else:
                raise ValueError('Error! Invalid validation measure.')

            current_epoch += 1
            if current_epoch % 1 == 0:
                perf_per_epoch.append(val_loss_mean)

                if current_epoch < config.max_keep + 1:
                    max_current_step.append(current_step)
                    max_perf_per_epoch.append(v_measure)

                    saver.save(sess=sess, save_path=os.path.join(log_path, 'model.ckpt'), global_step=current_step)
                    df_csv.loc[current_step, 'WEIGHT_PATH'] = \
                        os.path.join(log_path, 'model.ckpt-'+str(current_step))
                    df_csv.loc[current_step, 'AUC'] = val_auc
                    df_csv.loc[current_step, 'LOSS'] = val_loss_mean

                elif v_criteria:
                    if current_epoch >= config.save_epoch:
                        df_csv = df_csv.drop(max_current_step[0])
                        max_current_step.pop(0)
                        max_current_step.append(current_step)
                        max_perf_per_epoch.pop(0)
                        max_perf_per_epoch.append(v_measure)

                        saver.save(sess=sess, save_path=os.path.join(log_path, 'model.ckpt'), global_step=current_step)
                        df_csv.loc[current_step, 'WEIGHT_PATH'] = \
                            os.path.join(log_path, 'model.ckpt-'+str(current_step))
                        df_csv.loc[current_step, 'AUC'] = val_auc
                        df_csv.loc[current_step, 'LOSS'] = val_loss_mean

                df_csv.to_csv(os.path.join(result_path, result_name))

                if current_epoch == config.num_epoch:
                    break

        print('Training Complete...\n')
        sess.close()

    except KeyboardInterrupt:
        print('Result saved')
        df_csv.to_csv(os.path.join(result_path, result_name))


def validation():
    sess_config = tf.ConfigProto(log_device_placement=False, allow_soft_placement=True)
    sess_config.gpu_options.allow_growth = True

    saver = tf.train.Saver()
    num_examples = len(data_set.id_val)

    ckpt = tf.train.get_checkpoint_state(log_path)
    if not ckpt:
        raise ValueError('No checkpoint found in ' + log_path)

    weight_auc_path = os.path.join(config.data_path, config.exp_name, config.model_name,
                                   'result-%03d' % config.trial_serial)
    weight_auc_csv = pd.read_csv(os.path.join(
        weight_auc_path, '_'.join([config.exp_name, config.model_name, '%03d' % config.trial_serial])+'.csv'))

    if config.pre_weight == 'auc':
        weight_auc_csv = weight_auc_csv.sort_values('AUC', ascending=False)
    elif config.pre_weight == 'loss':
        weight_auc_csv = weight_auc_csv.sort_values('LOSS', ascending=True)
    else:
        raise ValueError('Error! Invalid pre weight criteria.')

    all_ckpt_paths = list(weight_auc_csv['WEIGHT_PATH'][config.select_weight:config.select_weight+1])

    num_ckpt = len(all_ckpt_paths)
    print('num_ckpt: ', num_ckpt)

    imgs = np.zeros([num_examples, model.img_h, model.img_w, model.img_c])
    cams = np.zeros([num_ckpt, num_examples, model.img_h, model.img_w, model.img_c])

    lbls = np.zeros([num_examples, ], dtype=np.int32)
    lgts = np.zeros([num_ckpt, num_examples, 1])
    probs = np.zeros([num_ckpt, num_examples, 1])

    val_x, val_y = None, None
    for ckpt_idx, ckpt_path in enumerate(all_ckpt_paths):
        print('Restoring: ' + ckpt_path)

        sess = tf.Session(config=sess_config)
        saver.restore(sess, ckpt_path)

        sess.run(data_set.val.init_op)
        val_x, val_y = [], []

        num_iter = int(np.ceil(float(num_examples) / config.batch_size))
        step = 0

        while step < num_iter:
            sys.stdout.write('Evaluation [{0}/{1}]\r'.format(len(val_y) // config.batch_size + 1,
                             -(-data_set.val.data_length // config.batch_size)))

            img, lbl, name, age, tm, dm, vas, _ = sess.run(data_set.val.next_batch)

            feed_dict = {model.images: img, model.labels: lbl, model.is_training: False}

            val_loss, val_lgt, val_prob, val_acc = \
                sess.run([model.loss, model.logits, model.prob, model.accuracy], feed_dict=feed_dict)

            if np.ndim(val_prob) == 1:
                val_prob = np.expand_dims(val_prob, axis=-1)

            val_x.extend(val_prob)
            val_y.extend(lbl)

            cam = sess.run(model.local, feed_dict=feed_dict)

            cams[ckpt_idx, step * config.batch_size:step * config.batch_size + len(lbl)] = cam
            probs[ckpt_idx, step * config.batch_size:step * config.batch_size + len(lbl)] = val_prob
            lgts[ckpt_idx, step * config.batch_size:step * config.batch_size + len(lbl)] = val_lgt

            if ckpt_idx == 0:
                imgs[step * config.batch_size:step * config.batch_size + len(lbl)] = img
                lbls[step * config.batch_size:step * config.batch_size + len(lbl)] = lbl

            step += 1

        sess.close()

    probs, lgts, cams = np.mean(probs, axis=0), np.mean(lgts, axis=0), np.mean(cams, axis=0)

    id_test = data_set.id_val

    prob_1, lgts_1 = np.squeeze(np.array(probs)), np.squeeze(np.array(lgts))

    result_csv = pd.DataFrame({'NUMBER': id_test, 'PROB': prob_1, 'LOGIT': lgts_1, 'LABEL': np.array(lbls)})
    result_name = '_'.join([config.model_name, config.npy_name, trial_serial_str, '%03d' % config.num_weight])+'.csv'
                             #config.pre_weight, '%03d' % config.select_weight

    result_csv.to_csv(os.path.join(result_path, result_name), index=False)

    fpr, tpr, _ = sklearn.metrics.roc_curve(val_y, prob_1, drop_intermediate=False)
    val_auc = sklearn.metrics.auc(fpr, tpr)

    print('Validation AUC: ', val_auc)
    print('Validation Complete...\n')

    if True:
        prs = pptx.Presentation()
        prs.slide_width, prs.slide_height = Inches(8*2), Inches(5*2)

        plt_batch = 20
        plt_step = 0
        plt_iter, plt_examples = int(np.ceil(num_examples / plt_batch)), num_examples

        while plt_step < plt_iter:

            if plt_examples >= plt_batch:
                len_batch = plt_batch
            else:
                len_batch = plt_examples

            images_batch = imgs[plt_step * plt_batch:plt_step * plt_batch + len_batch]
            labels_batch = lbls[plt_step * plt_batch:plt_step * plt_batch + len_batch]
            names_batch = id_test[plt_step * plt_batch:plt_step * plt_batch + len_batch]

            probs_batch = probs[plt_step * plt_batch:plt_step * plt_batch + len_batch]
            cams_batch = cams[plt_step * plt_batch:plt_step * plt_batch + len_batch]

            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)

            show_cam(cams_batch, probs_batch, images_batch, labels_batch, names_batch, 'LABEL')
            fig_name = '_'.join([config.exp_name, config.model_name, config.npy_name, trial_serial_str,
                                 '%03d' % plt_step]) + '.png'
            fig_path = os.path.join(cam_path, fig_name)
            plt.savefig(fig_path, bbox_inches='tight')
            slide.shapes.add_picture(fig_path, Inches(0), Inches(0), width=Inches(8 * 2))
            os.remove(fig_path)
            plt_step += 1
            plt_examples -= plt_batch

        print('plt_examples check: ', plt_examples)
        ppt_name = os.path.join(ppt_path, '_'.join([config.exp_name, config.model_name, config.npy_name,
                                                    trial_serial_str, '%03d' % config.num_weight,
                                                    config.pre_weight, '%03d' % config.select_weight]) + '.pptx')
        prs.save(ppt_name)
        print('Saved: ', ppt_name)


def show_cam(cams, probs, images, labels, names, side_label, num_rows=5, num_cols=8, fig_size=(8*2, 5*2)):
    batch_size = cams.shape[0]
    fig, ax = plt.subplots(num_rows, num_cols, figsize=fig_size)
    axoff_fun = np.vectorize(lambda ax: ax.axis('off'))
    axoff_fun(ax)

    for i in range(batch_size):
        prob = '%.2f' % probs[i]
        lbl = int(labels[i])
        show_image = np.squeeze(images[i])
        cam = np.squeeze(cams[i])
        img_row, img_col = int(i % num_rows), int(i / num_rows) * 2

        ori_title = ' '.join([names[i], side_label + ': '+str(lbl)])
        cam_title = side_label+' Pred: '+str(prob)

        ax[img_row, img_col].imshow(show_image, cmap='bone')
        ax[img_row, img_col+1].imshow(show_image, cmap='bone')
        ax[img_row, img_col+1].imshow(cam, cmap=plt.cm.jet, alpha=0.5, interpolation='nearest')

        if (lbl == 0 and probs[i] < 0.5) or (lbl == 1 and probs[i] >= 0.5):
            txt_color = 'blue'
        else:
            txt_color = 'red'
        ax[img_row, img_col].set_title(ori_title, fontsize=7, color=txt_color)
        ax[img_row, img_col+1].set_title(cam_title, fontsize=7, color=txt_color)


if __name__ == '__main__':
    if config.train:
        print('Training')
        training()
    else:
        print('Validation')
        validation()